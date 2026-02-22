#!/usr/bin/env python3
"""Generate a hackathon pitch deck for MakeItSo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DAVIDSON_RED = RGBColor(0x8C, 0x1D, 0x1D)
DAVIDSON_DARK = RGBColor(0x7A, 0x19, 0x19)
NAVY = RGBColor(0x1E, 0x2A, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFB)
MEDIUM_GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK_TEXT = RGBColor(0x11, 0x11, 0x11)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        solidFill = shape.fill._fill
        srgbClr = solidFill.find(f'{{{ns}}}solidFill/{{{ns}}}srgbClr')
        if srgbClr is not None:
            alpha_elem = etree.SubElement(srgbClr, f'{{{ns}}}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=18,
                       default_color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                       font_name="Calibri"):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# Accent bar at top
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), DAVIDSON_RED)

# Decorative circles (subtle)
add_circle(slide, Inches(10.5), Inches(0.5), Inches(3), RGBColor(0x2A, 0x3A, 0x4E))
add_circle(slide, Inches(-1), Inches(5), Inches(2.5), RGBColor(0x2A, 0x3A, 0x4E))

# Sparkle icon placeholder - red rounded rect
icon_shape = add_rounded_rect(slide, Inches(5.9), Inches(1.6), Inches(1.5), Inches(1.5), DAVIDSON_RED)
add_text_box(slide, Inches(5.9), Inches(1.72), Inches(1.5), Inches(1.3),
             "MIS", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Georgia")

# Title
add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10.3), Inches(1.2),
             "MakeItSo", font_size=60, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Georgia")

# Tagline
add_text_box(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.8),
             "AI-powered course planning that connects what you study to where you're going.",
             font_size=22, color=RGBColor(0xA0, 0xAE, 0xC0), bold=False,
             alignment=PP_ALIGN.CENTER)

# Bottom line
add_text_box(slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.5),
             "hack@DAVIDSON 2025  |  Built for Davidson College Students",
             font_size=14, color=RGBColor(0x6B, 0x7B, 0x90), bold=False,
             alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, DAVIDSON_RED)

# Section label
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "THE PROBLEM", font_size=13, color=DAVIDSON_RED, bold=True,
             font_name="Calibri")

# Title
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1),
             "Students pick courses blind to career outcomes.",
             font_size=40, color=NAVY, bold=True, font_name="Georgia")

problems = [
    ("No career connection", "Students choose courses with zero visibility into how they map to actual careers."),
    ("Information overload", "463 courses across 24 departments. Generic advisor meetings aren't enough."),
    ("Networking anxiety", "Alumni exist who can help, but students don't know what to say or who to contact."),
    ("Professor roulette", "RateMyProfessors is noisy. There's no synthesized, actionable professor intel."),
]

y_start = Inches(2.5)
for i, (title, desc) in enumerate(problems):
    x = Inches(0.8) + Inches(3.05) * i
    card = add_rounded_rect(slide, x, y_start, Inches(2.8), Inches(3.6), LIGHT_GRAY)

    # Number circle
    num_circle = add_circle(slide, x + Inches(0.2), y_start + Inches(0.25), Inches(0.55), DAVIDSON_RED)
    add_text_box(slide, x + Inches(0.2), y_start + Inches(0.3), Inches(0.55), Inches(0.45),
                 str(i + 1), font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.2), y_start + Inches(1.05), Inches(2.4), Inches(0.5),
                 title, font_size=18, color=NAVY, bold=True, font_name="Georgia")

    add_text_box(slide, x + Inches(0.2), y_start + Inches(1.6), Inches(2.4), Inches(1.8),
                 desc, font_size=14, color=MEDIUM_GRAY, line_spacing=1.4)


# ============================================================
# SLIDE 3: THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "THE SOLUTION", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1),
             'Don\'t ask "what courses?" Ask "what career?"',
             font_size=40, color=NAVY, bold=True, font_name="Georgia")

add_text_box(slide, Inches(0.8), Inches(1.85), Inches(10), Inches(0.7),
             "MakeItSo flips course planning on its head. Start with your dream career, and AI builds your entire academic roadmap backward\u2014courses, people, activities, and timeline.",
             font_size=17, color=MEDIUM_GRAY, line_spacing=1.5)

# Two pathways
path_data = [
    ("Plan Forward", "Start with your interests", ACCENT_GREEN,
     ["Tell us what fascinates you", "Get AI course recommendations", "Explore 24 career paths", "Discover new possibilities"]),
    ("Plan Backward", "Start with your dream career", DAVIDSON_RED,
     ["Pick your target role", "AI generates full action plan", "Courses + People + Activities", "Follow your roadmap to success"]),
]

for i, (title, subtitle, color, steps) in enumerate(path_data):
    x = Inches(1.5) + Inches(5.5) * i
    y = Inches(3.1)

    card = add_rounded_rect(slide, x, y, Inches(4.8), Inches(3.8), LIGHT_GRAY)

    # Color bar at top of card
    add_rect(slide, x, y, Inches(4.8), Inches(0.06), color)

    add_text_box(slide, x + Inches(0.4), y + Inches(0.35), Inches(4), Inches(0.45),
                 title, font_size=24, color=NAVY, bold=True, font_name="Georgia")
    add_text_box(slide, x + Inches(0.4), y + Inches(0.85), Inches(4), Inches(0.35),
                 subtitle, font_size=14, color=color, bold=True)

    for j, step in enumerate(steps):
        step_y = y + Inches(1.5) + Inches(0.5) * j
        num_c = add_circle(slide, x + Inches(0.4), step_y, Inches(0.3), color)
        add_text_box(slide, x + Inches(0.4), step_y + Inches(0.02), Inches(0.3), Inches(0.26),
                     str(j + 1), font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.9), step_y + Inches(0.02), Inches(3.5), Inches(0.3),
                     step, font_size=14, color=DARK_TEXT)


# ============================================================
# SLIDE 4: KEY FEATURES (1 of 2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "KEY FEATURES", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             "Six AI-powered tools, one platform.",
             font_size=40, color=NAVY, bold=True, font_name="Georgia")

features = [
    ("Smart Course Explorer", "463 real Davidson courses with live enrollment, RateMyProfessors ratings, and AI-generated insights showing key topics, skills gained, and career connections.",
     ACCENT_BLUE),
    ("AI Career Planner", "Tell us your dream career. Gemini generates a complete action plan: recommended courses, people to meet, activities to do, and Davidson-specific insights.",
     DAVIDSON_RED),
    ("4-Year Roadmap Builder", "AI creates optimized semester-by-semester schedules respecting prerequisites, distribution requirements, and workload balance. Adjustable specificity levels 1\u20135.",
     ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(features):
    x = Inches(0.8) + Inches(4.1) * i
    y = Inches(2.2)
    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.5), LIGHT_GRAY)
    add_rect(slide, x, y, Inches(3.8), Inches(0.06), color)

    # Icon circle
    icon_c = add_circle(slide, x + Inches(0.3), y + Inches(0.4), Inches(0.6), color)
    icons = ["\u2605", "\u25B6", "\u2726"]
    add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(0.6), Inches(0.5),
                 icons[i], font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.5),
                 title, font_size=19, color=NAVY, bold=True, font_name="Georgia")

    add_text_box(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.2), Inches(2.5),
                 desc, font_size=13, color=MEDIUM_GRAY, line_spacing=1.5)


# ============================================================
# SLIDE 5: KEY FEATURES (2 of 2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "KEY FEATURES", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             "Real data. Real people. Real outcomes.",
             font_size=40, color=NAVY, bold=True, font_name="Georgia")

features2 = [
    ("Professor Insights", "AI synthesizes RateMyProfessors reviews into balanced summaries with strengths, considerations, and a specific tip for success\u2014actionable, not discouraging.",
     ACCENT_PURPLE),
    ("Alumni Network", "40+ curated Davidson alumni across McKinsey, Goldman Sachs, Google, Microsoft, and more. Browse by career field, company, and graduation year.",
     ACCENT_AMBER),
    ("Cold Email Generator", "Solves networking anxiety. AI writes personalized, professional outreach emails tailored to each alumnus's role and the student's interests.",
     RGBColor(0xEC, 0x48, 0x99)),
]

for i, (title, desc, color) in enumerate(features2):
    x = Inches(0.8) + Inches(4.1) * i
    y = Inches(2.2)
    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.5), LIGHT_GRAY)
    add_rect(slide, x, y, Inches(3.8), Inches(0.06), color)

    icon_c = add_circle(slide, x + Inches(0.3), y + Inches(0.4), Inches(0.6), color)
    icons2 = ["\u2606", "\u25CF", "\u2709"]
    add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(0.6), Inches(0.5),
                 icons2[i], font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.5),
                 title, font_size=19, color=NAVY, bold=True, font_name="Georgia")

    add_text_box(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.2), Inches(2.5),
                 desc, font_size=13, color=MEDIUM_GRAY, line_spacing=1.5)


# ============================================================
# SLIDE 6: HOW IT WORKS (Architecture)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
             "UNDER THE HOOD", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
             "Built with a modern, production-grade stack.",
             font_size=36, color=WHITE, bold=True, font_name="Georgia")

# Tech stack boxes
stack_items = [
    ("Frontend", "Next.js 14 + TypeScript\nshadcn/ui + Tailwind CSS\nFramer Motion animations\nRecharts visualizations", ACCENT_BLUE),
    ("Backend", "Next.js API Routes (serverless)\nMongoDB Atlas + Mongoose\nNextAuth.js authentication\nbcrypt password hashing", ACCENT_GREEN),
    ("AI Engine", "Google Gemini 3 Flash\nStructured JSON outputs\nSmart caching (80% hit rate)\n6 distinct AI features", DAVIDSON_RED),
    ("Data Sources", "Davidson College API (live)\nRateMyProfessors GraphQL\n463 courses indexed\n40+ alumni curated", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(stack_items):
    x = Inches(0.6) + Inches(3.15) * i
    y = Inches(2.3)
    card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(3.6), RGBColor(0x2A, 0x3A, 0x4E))
    add_rect(slide, x, y, Inches(2.9), Inches(0.06), color)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.3), Inches(0.4),
                 title, font_size=20, color=color, bold=True, font_name="Georgia")

    # Description lines
    for j, line in enumerate(desc.split("\n")):
        add_text_box(slide, x + Inches(0.3), y + Inches(1.0) + Inches(0.5) * j,
                     Inches(2.3), Inches(0.4),
                     line, font_size=13, color=RGBColor(0xA0, 0xAE, 0xC0))

# Deployment note
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
             "Deployed on Vercel  |  MongoDB Atlas  |  Serverless, auto-scaling  |  <2s page loads",
             font_size=13, color=RGBColor(0x6B, 0x7B, 0x90), alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: BY THE NUMBERS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "BY THE NUMBERS", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             "Real data, real scale, real impact.",
             font_size=40, color=NAVY, bold=True, font_name="Georgia")

stats = [
    ("463", "Davidson courses\nfully indexed", ACCENT_BLUE),
    ("24", "Career paths\nmapped", DAVIDSON_RED),
    ("40+", "Alumni contacts\ncurated", ACCENT_GREEN),
    ("6", "AI-powered\nfeatures", ACCENT_PURPLE),
    ("128", "Credits tracked\nto graduation", ACCENT_AMBER),
    ("22", "API endpoints\nbuilt", RGBColor(0xEC, 0x48, 0x99)),
]

for i, (number, label, color) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = Inches(1.2) + Inches(3.8) * col
    y = Inches(2.4) + Inches(2.5) * row

    card = add_rounded_rect(slide, x, y, Inches(3.2), Inches(2.0), LIGHT_GRAY)

    add_text_box(slide, x, y + Inches(0.25), Inches(3.2), Inches(0.8),
                 number, font_size=48, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Georgia")

    for j, line in enumerate(label.split("\n")):
        add_text_box(slide, x, y + Inches(1.1) + Inches(0.3) * j, Inches(3.2), Inches(0.3),
                     line, font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8: WHAT MAKES US DIFFERENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "WHY WE WIN", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             "This isn't another generic course planner.",
             font_size=40, color=NAVY, bold=True, font_name="Georgia")

differentiators = [
    ("Backward-from-career planning", "We flip the paradigm. Most tools ask \u201cwhat courses should I take?\u201d We ask \u201cwhat career do you want?\u201d and build backward."),
    ("Davidson-native depth", "Not a generic tool. 463 real Davidson courses, real faculty, real alumni, live enrollment from Davidson\u2019s own API."),
    ("AI-synthesized professor reviews", "Raw RMP is noisy. Our AI reads actual reviews and creates balanced, actionable summaries with tips for success."),
    ("Career intelligence on every course", "Every course shows career relevance: \u201cCSC 121 is 90% relevant to Software Engineering, 45% to Data Science.\u201d"),
    ("Alumni network + cold email outreach", "We don\u2019t just show alumni\u2014we generate personalized outreach emails so students actually network."),
]

for i, (title, desc) in enumerate(differentiators):
    y = Inches(2.2) + Inches(0.95) * i
    # Accent dot
    add_circle(slide, Inches(0.9), y + Inches(0.1), Inches(0.15), DAVIDSON_RED)
    add_text_box(slide, Inches(1.3), y, Inches(3.5), Inches(0.4),
                 title, font_size=17, color=NAVY, bold=True, font_name="Georgia")
    add_text_box(slide, Inches(1.3), y + Inches(0.35), Inches(10.5), Inches(0.5),
                 desc, font_size=13, color=MEDIUM_GRAY, line_spacing=1.4)


# ============================================================
# SLIDE 9: DEMO FLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
             "DEMO WALKTHROUGH", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
             "From sign-up to career roadmap in 60 seconds.",
             font_size=36, color=WHITE, bold=True, font_name="Georgia")

demo_steps = [
    ("1", "Sign Up", "Create account,\nset major & interests", ACCENT_BLUE),
    ("2", "Explore Careers", "Browse 24 paths,\npick your target", DAVIDSON_RED),
    ("3", "Generate Plan", "AI builds courses,\npeople, activities", ACCENT_GREEN),
    ("4", "Build Schedule", "Drag & drop into\n4-year roadmap", ACCENT_PURPLE),
    ("5", "Network", "Find alumni, generate\ncold emails with AI", ACCENT_AMBER),
]

for i, (num, title, desc, color) in enumerate(demo_steps):
    x = Inches(0.65) + Inches(2.5) * i
    y = Inches(2.3)

    card = add_rounded_rect(slide, x, y, Inches(2.2), Inches(3.5), RGBColor(0x2A, 0x3A, 0x4E))

    # Number circle
    nc = add_circle(slide, x + Inches(0.75), y + Inches(0.35), Inches(0.7), color)
    add_text_box(slide, x + Inches(0.75), y + Inches(0.42), Inches(0.7), Inches(0.55),
                 num, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name="Georgia")

    add_text_box(slide, x + Inches(0.15), y + Inches(1.35), Inches(1.9), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Georgia")

    for j, line in enumerate(desc.split("\n")):
        add_text_box(slide, x + Inches(0.15), y + Inches(1.9) + Inches(0.35) * j,
                     Inches(1.9), Inches(0.3),
                     line, font_size=13, color=RGBColor(0xA0, 0xAE, 0xC0),
                     alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        add_text_box(slide, x + Inches(2.2), y + Inches(1.3), Inches(0.3), Inches(0.4),
                     "\u25B8", font_size=20, color=RGBColor(0x4A, 0x5A, 0x6E),
                     alignment=PP_ALIGN.CENTER)

# Note at bottom
add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
             "Live demo available \u2014 try it yourself at the booth!",
             font_size=15, color=RGBColor(0x6B, 0x7B, 0x90), alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: WHAT'S NEXT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, DAVIDSON_RED)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "WHAT'S NEXT", font_size=13, color=DAVIDSON_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             "From Davidson prototype to institutional platform.",
             font_size=40, color=NAVY, bold=True, font_name="Georgia")

future_items = [
    ("NOW", "hack@DAVIDSON MVP",
     "463 courses, 24 career paths, 6 AI features, live enrollment sync, alumni network, cold email gen.",
     ACCENT_GREEN),
    ("NEXT", "Campus-Wide Launch",
     "Student beta testing, advisor integrations, mobile optimization, feedback loop with academic deans.",
     ACCENT_BLUE),
    ("FUTURE", "Multi-School Platform",
     "Adapt the engine for other liberal arts colleges. B2B licensing model for institutional career offices.",
     ACCENT_PURPLE),
]

for i, (phase, title, desc, color) in enumerate(future_items):
    x = Inches(0.8) + Inches(4.1) * i
    y = Inches(2.3)
    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.2), LIGHT_GRAY)
    add_rect(slide, x, y, Inches(3.8), Inches(0.06), color)

    # Phase badge
    badge = add_rounded_rect(slide, x + Inches(0.3), y + Inches(0.4), Inches(1.2), Inches(0.4), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.42), Inches(1.2), Inches(0.35),
                 phase, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.2), Inches(0.5),
                 title, font_size=19, color=NAVY, bold=True, font_name="Georgia")

    add_text_box(slide, x + Inches(0.3), y + Inches(1.7), Inches(3.2), Inches(2.2),
                 desc, font_size=13, color=MEDIUM_GRAY, line_spacing=1.5)


# ============================================================
# SLIDE 11: CLOSING / CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), DAVIDSON_RED)

# Decorative circles
add_circle(slide, Inches(10.5), Inches(0.5), Inches(3), RGBColor(0x2A, 0x3A, 0x4E))
add_circle(slide, Inches(-1), Inches(5), Inches(2.5), RGBColor(0x2A, 0x3A, 0x4E))

# Icon
icon_shape = add_rounded_rect(slide, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5), DAVIDSON_RED)
add_text_box(slide, Inches(5.9), Inches(1.32), Inches(1.5), Inches(1.3),
             "MIS", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Georgia")

add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1),
             "Your degree. Your career. One plan.",
             font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Georgia")

add_text_box(slide, Inches(2), Inches(4.1), Inches(9.3), Inches(0.7),
             "MakeItSo gives every Davidson student an AI career advisor\u2014free, personalized, and always available.",
             font_size=18, color=RGBColor(0xA0, 0xAE, 0xC0),
             alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# CTA button shape
cta = add_rounded_rect(slide, Inches(5.1), Inches(5.2), Inches(3.1), Inches(0.65), DAVIDSON_RED)
add_text_box(slide, Inches(5.1), Inches(5.25), Inches(3.1), Inches(0.55),
             "Try the Live Demo", font_size=18, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.3), Inches(9.3), Inches(0.5),
             "hack@DAVIDSON 2025  |  Questions? Come talk to us!",
             font_size=14, color=RGBColor(0x6B, 0x7B, 0x90),
             alignment=PP_ALIGN.CENTER)


# Save
output_path = "/home/user/MakeItSo/MakeItSo_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
